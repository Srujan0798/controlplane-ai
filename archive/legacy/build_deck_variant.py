#!/usr/bin/env python3
"""ControlPlane.ai — Variant Rebuild. Generates a second, comparable version
of slides 2-4 on top of the already-built ControlPlane_ControlPlane-ai.pptx
(team slide, slide order and file-size compression are inherited from it)."""
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ================================================================== PALETTE
PURPLE      = C(0xA0, 0x00, 0xFF)
DEEP        = C(0x46, 0x00, 0x73)
INK         = C(0x1A, 0x1A, 0x1A)
INK2        = C(0x5A, 0x5A, 0x60)
INK3        = C(0x94, 0x94, 0x9C)
HAIR        = C(0xE2, 0xE2, 0xE6)
SOFT        = C(0xF4, 0xF3, 0xF7)
WHITE       = C(0xFF, 0xFF, 0xFF)
DARK        = C(0x16, 0x18, 0x1C)
RED         = C(0xD1, 0x34, 0x38)
RED_SOFT    = C(0xF6, 0xDD, 0xDD)
BLUE        = C(0x22, 0x4B, 0xFF)
BLUE_SOFT   = C(0xDF, 0xE4, 0xFB)
AMBER       = C(0xB2, 0x6A, 0x00)
AMBER_SOFT  = C(0xFA, 0xEE, 0xD9)
TEAL        = C(0x1B, 0x63, 0x50)   # was missing in the source script — cost axis colour
GREY_SOFT   = C(0xF1, 0xF1, 0xF3)
GREY_MID    = C(0xC6, 0xC6, 0xCE)
F = "Arial"
FM = "Menlo"  # Consolas is not installed here; Menlo is the macOS mono fallback

W, H = 13.33, 7.50
ML, MR = 0.90, 0.90
CW = W - ML - MR

# ================================================================== HELPERS
def txt(sl, x, y, w, h, s, size=13, bold=False, color=INK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=F, line=None, italic=False):
    tb = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, l in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line: p.line_spacing = line
        r = p.add_run(); r.text = l
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb

def box(sl, x, y, w, h, fill=None, linec=None, lw=1.0):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if linec is None: s.line.fill.background()
    else: s.line.color.rgb = linec; s.line.width = Pt(lw)
    s.shadow.inherit = False; return s

def ln(sl, x1, y1, x2, y2, color=HAIR, w=1.0):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); return c

def arrow(sl, x, y, length=0.44, color=C(0xC4,0xC4,0xCC)):
    ln(sl, x, y, x+length, y, color=color, w=1.5)
    ln(sl, x+length-0.11, y-0.06, x+length, y, color=color, w=1.5)
    ln(sl, x+length-0.11, y+0.06, x+length, y, color=color, w=1.5)

# ================================================================== LOAD
SRC = "ControlPlane_ControlPlane-ai.pptx"
OUT = "ControlPlane_ControlPlane-ai_VARIANT.pptx"
prs = Presentation(SRC)

def clear(sl):
    for sh in list(sl.shapes):
        sh._element.getparent().remove(sh._element)

# ================================================================== SLIDE 2: PROBLEM
s2 = prs.slides[2]; clear(s2)
txt(s2, ML, 0.48, CW, 0.22, "THE PROBLEM", size=10, bold=True, color=PURPLE)
txt(s2, ML, 0.72, CW, 1.05, "Every filter passed.\nThe money moved anyway.",
    size=36, bold=True, color=DEEP, line=1.05)

box(s2, ML, 2.12, CW, 1.42, fill=DARK)
txt(s2, ML+0.45, 2.38, CW-0.9, 0.36,
    "Approved.  Refund of ₹1,84,000 issued under clause 7.2",
    font=FM, size=18, color=C(0xF2,0xF2,0xEF))
txt(s2, ML+0.45, 2.82, CW-0.9, 0.3, "clause 7.2 does not exist",
    font=FM, size=15, bold=True, color=RED)

facts = [("EXECUTED","Tuesday"), ("DISCOVERED","Friday"), ("MODEL CONFIDENCE","0.94")]
for i,(k,v) in enumerate(facts):
    x = ML + i*(CW/3)
    txt(s2, x, 3.78, CW/3-0.5, 0.2, k, size=9, bold=True, color=INK3)
    txt(s2, x, 4.00, CW/3-0.5, 0.4, v, size=24, bold=True, color=INK)

ln(s2, ML, 4.62, ML+CW, 4.62, color=HAIR)
txt(s2, ML, 4.82, CW, 0.75,
    "The system didn't fail.\nIt was never asked to prove anything.",
    size=28, bold=True, color=DEEP, line=1.05)

# ================================================================== SLIDE 3: SOLUTION
s3 = prs.slides[3]; clear(s3)
txt(s3, ML, 0.48, CW, 0.25,
    "It used to be a bad paragraph. It is now an executed transaction.",
    size=11, italic=True, color=INK3)
txt(s3, ML, 0.78, CW, 0.25, "PROPOSED SOLUTION", size=10, bold=True, color=PURPLE)
txt(s3, ML, 0.98, CW, 0.90,
    "An AI response is a set of claims\nrequesting permission to act.",
    size=34, bold=True, color=DEEP, line=1.05)

NW, NH, GAP = 2.35, 1.50, 0.58
gy = 2.48
NX = [ML + i*(NW+GAP) for i in range(4)]

box(s3, NX[0], gy, NW, NH, fill=SOFT)
txt(s3, NX[0]+0.22, gy+0.22, NW-0.44, 0.4, "STEP", size=22, bold=True, color=DEEP)
txt(s3, NX[0]+0.22, gy+0.68, NW-0.44, 0.5, "tool calls\nand retrieval", size=11, color=INK2, line=1.3)
box(s3, NX[1], gy, NW, NH, fill=SOFT)
txt(s3, NX[1]+0.22, gy+0.22, NW-0.44, 0.4, "SPAN", size=22, bold=True, color=DEEP)
txt(s3, NX[1]+0.22, gy+0.68, NW-0.44, 0.5, "source · ACL · hash", size=11, color=INK2, line=1.3)
box(s3, NX[2], gy, NW, NH, fill=SOFT)
txt(s3, NX[2]+0.22, gy+0.22, NW-0.44, 0.4, "CLAIM", size=22, bold=True, color=DEEP)
txt(s3, NX[2]+0.22, gy+0.60, NW-0.44, 0.2, "typed, atomic", size=11, color=INK2)
for i in range(6):
    fill = RED if i==5 else GREY_MID
    box(s3, NX[2]+0.22+i*0.30, gy+0.88, 0.24, 0.16, fill=fill)
txt(s3, NX[2]+0.22, gy+1.10, NW-0.44, 0.2, "5 bind.  1 does not.", size=10, bold=True, color=RED)
box(s3, NX[3], gy, NW, NH, fill=WHITE, linec=INK, lw=2.0)
txt(s3, NX[3]+0.22, gy+0.22, NW-0.44, 0.4, "ACTION", size=22, bold=True, color=INK)
txt(s3, NX[3]+0.22, gy+0.68, NW-0.44, 0.5, "refund ₹1,84,000", size=11, color=INK2, line=1.3)

for i in range(3): arrow(s3, NX[i]+NW+0.07, gy+NH/2)
bx = NX[1]-0.18
for k in range(5):
    y1 = gy-0.14+k*0.07; y2 = y1+0.045
    ln(s3, bx, y1, bx, y2, color=PURPLE, w=2.0)
txt(s3, bx+0.10, gy-0.22, 3.8, 0.18,
    "CAPTURED HERE — OUTSIDE THE MODEL", size=9, bold=True, color=PURPLE)
gx = NX[2]+NW+0.29; gy2 = gy+NH/2
ln(s3, gx-0.12, gy2-0.14, gx+0.12, gy2+0.14, color=RED, w=3.0)
ln(s3, gx+0.12, gy2-0.14, gx-0.12, gy2+0.14, color=RED, w=3.0)

by = gy+NH+0.42
leaders = [("COST","4 of 9 steps grounded nothing",TEAL),
           ("RESPONSIBILITY","1 span the caller may not read",AMBER),
           ("PERFORMANCE","1 claim with no span",BLUE)]
for i,(t,d,col) in enumerate(leaders):
    x = ML + i*(CW/3)
    box(s3, x, by, 0.06, 0.38, fill=col)
    txt(s3, x+0.14, by, CW/3-0.4, 0.2, t, size=9, bold=True, color=col)
    txt(s3, x+0.14, by+0.20, CW/3-0.4, 0.3, d, size=14, bold=True, color=INK)

txt(s3, ML, by+0.65, CW, 0.4, "Three dimensions. One graph.", size=26, bold=True, color=DEEP)

# ================================================================== SLIDE 4: DIFFERENT
s4 = prs.slides[4]; clear(s4)
txt(s4, ML, 0.42, CW, 0.22, "WHY THIS IS DIFFERENT", size=10, bold=True, color=PURPLE)
txt(s4, ML, 0.66, CW, 0.85,
    "The same unproven claim annotates a draft\nand holds a payment.",
    size=32, bold=True, color=DEEP, line=1.05)

my = 1.72; c0w = 1.65; cw = (CW-c0w)/4; rh = 0.58; hh = 0.26
headers = ["Contradicted","Unsupported","Hedged","Unknown"]
for j,h in enumerate(headers):
    x = ML+c0w+j*cw
    txt(s4, x, my, cw-0.06, hh, h, size=9, bold=True, color=INK3, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

data = [("R3","irreversible",["BLOCK","ESCALATE","ESCALATE","ESCALATE"]),
        ("R2","reversible write",["BLOCK","EDIT","EDIT","ESCALATE"]),
        ("R1","read-only",["EDIT","EDIT","PASS","PASS"]),
        ("R0","internal draft",["PASS","PASS","PASS","PASS"])]
SHADE = {"BLOCK":(RED_SOFT,RED), "ESCALATE":(BLUE_SOFT,BLUE),
         "EDIT":(AMBER_SOFT,AMBER), "PASS":(GREY_SOFT,C(0x8A,0x8A,0x92))}

for i,(tier,desc,cells) in enumerate(data):
    y = my+hh+0.08+i*(rh+0.06)
    txt(s4, ML, y+0.10, c0w-0.15, 0.35, tier, size=17, bold=True, color=INK)
    txt(s4, ML, y+0.42, c0w-0.15, 0.2, desc, size=9, color=INK3)
    for j,cell in enumerate(cells):
        x = ML+c0w+j*cw; bg,fg = SHADE[cell]
        box(s4, x, y, cw-0.06, rh, fill=bg)
        txt(s4, x, y, cw-0.06, rh, cell, size=12, bold=True, color=fg,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

brx = ML+c0w+1*cw-0.03
bt = my+hh+0.08; bb = bt+4*(rh+0.06)-0.06
ln(s4, brx-0.06, bt+0.10, brx-0.06, bb-0.10, color=INK3, w=1.0)
ln(s4, brx-0.06, bt+0.10, brx, bt+0.10, color=INK3, w=1.0)
ln(s4, brx-0.06, bb-0.10, brx, bb-0.10, color=INK3, w=1.0)

sy = my+hh+4*(rh+0.06)+0.16
txt(s4, ML, sy, CW, 0.20,
    "One unproven claim, four outcomes.   Unsupported × R3  ·  clause 7.2 → ESCALATE · ₹1,84,000 held",
    size=11, bold=True, color=BLUE)
txt(s4, ML, sy+0.26, CW, 0.18,
    "EDIT  —  strip or re-ground the named span, never a rewrite.     ESCALATE  —  inline hold; ships claim + spans + verdict.",
    size=9, color=INK3)
txt(s4, ML, sy+0.46, CW, 0.18,
    "Bias  —  counterfactual flip rate, route-level, CI excludes zero.     Safety  —  typed interlocks: tool × args × irreversibility.",
    size=9, color=INK3)
txt(s4, ML, sy+0.66, CW, 0.18,
    "Hard gate on actions, not tokens. Text streams with a short hold-back.",
    size=9, bold=True, color=INK2)

txt(s4, ML, 6.12, CW, 0.26,
    "We publish our own per-route false-negative rate. Everyone else publishes precision.",
    size=13, bold=True, color=INK2)

box(s4, 0, 6.48, W, 1.02, fill=DEEP)
txt(s4, ML, 6.72, CW, 0.55, "Now nothing acts until it can prove it should.",
    size=26, bold=True, color=WHITE)

# ================================================================== SAVE
prs.save(OUT)
print(f"Saved: {OUT}")
