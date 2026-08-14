#!/usr/bin/env python3
"""Official AIC template package — Cover · Team · 3 content · Video · Thank You.

Content slides 3–5 and Video are FULL-BLEED posters (no letterbox chrome).
Team slide is native template with real data, no leftover instruction text.
Also writes form-only 3-slide twin via build_submission when posters exist.
"""
from __future__ import annotations

import os
import subprocess
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor as C
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I, Pt, Emu
from PIL import Image

ROOT = Path(__file__).resolve().parent
SRC = ROOT / "AIC_Talent-Brand_PPT-Template (1).pptx"
OUT = ROOT / "ControlPlane_ControlPlane-ai.pptx"
SK = Path(os.path.expanduser("~/.grok/bundled/skills/pptx/scripts"))
SW, SH = 13.333333, 7.5
MAX_EDGE = 1920

PURPLE = C(0x74, 0x00, 0xC0)
DEEP = C(0x45, 0x00, 0x73)
PINK = C(0xA1, 0x00, 0xFF)
INK = C(0x00, 0x00, 0x00)
MUTED = C(0x5A, 0x56, 0x60)
WHITE = C(0xFF, 0xFF, 0xFF)
PLATE = C(0x10, 0x11, 0x14)
RED = C(0xD0, 0x18, 0x20)
FONT = "Arial"


def shape_text(sh) -> str:
    if not sh.has_text_frame:
        return ""
    return " ".join(p.text for p in sh.text_frame.paragraphs).strip()


def set_run(run, text, size=None, bold=None, color=None, italic=None, name=FONT):
    run.text = text
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def wipe_tf(tf):
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    return p0


def put_text(shape, lines, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p0 = wipe_tf(tf)
    p0.alignment = align
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        set_run(r, line, size=size, bold=bold, color=color)


def add_tb(slide, x, y, w, h, lines, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p0 = wipe_tf(tf)
    p0.alignment = align
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        set_run(r, line, size=size, bold=bold, color=color, italic=italic)
    return tb


def clear_shapes(slide):
    for sh in list(slide.shapes):
        el = sh._element
        el.getparent().remove(el)


def _compress_png(path: Path) -> BytesIO:
    im = Image.open(path).convert("RGB")
    w, h = im.size
    scale = min(1.0, MAX_EDGE / max(w, h))
    if scale < 1.0:
        im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    buf = BytesIO()
    im.save(buf, "JPEG", quality=90, optimize=True, progressive=True)
    buf.seek(0)
    return buf


def full_bleed(slide, png: Path):
    """Replace entire slide with a full-bleed poster — no letterbox, no title chrome."""
    clear_shapes(slide)
    assert png.exists(), png
    slide.shapes.add_picture(_compress_png(png), Emu(0), Emu(0), I(SW), I(SH))


def fill_team(slide):
    """Official AIC Team details — two people, no leftover instruction text."""
    for sh in list(slide.shapes):
        y = sh.top.inches if sh.top else 0
        t = shape_text(sh)
        kill = False
        if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            kill = True
        if sh.name in ("Rectangle 29", "Straight Connector 63", "TextBox 18"):
            kill = True
        if y > 4.5 and (sh.name.startswith("Straight") or t in ("", " ", "Name", "Photo")):
            kill = True
        if "mandatory" in t.lower():
            kill = True
        if kill:
            sh._element.getparent().remove(sh._element)

    for sh in slide.shapes:
        t = shape_text(sh)
        if not sh.has_text_frame and not sh.has_table:
            continue
        x = sh.left.inches
        y = sh.top.inches
        if sh.has_table:
            try:
                sh.table.cell(0, 0).text = "TEAM NAME:"
                sh.table.cell(0, 1).text = "ControlPlane"
            except Exception:
                pass
            continue
        if "Team Leader" in t or (t.startswith("Name") and x < 6 and y < 3):
            put_text(sh, "Choda Srujan Sai  ·  Team Leader", size=18, bold=True, color=PINK)
        elif t == "Name" and x > 8:
            put_text(sh, "Dhrithika", size=18, bold=True, color=PINK)
        elif ("College" in t or t.startswith("College")) and x < 6 and y < 4.5:
            put_text(sh, [
                "College:  IIT Gandhinagar",
                "Stream:  CSE (Civil Engg. minor)",
                "Year of graduation:  2027",
            ], size=13, color=INK)
        elif ("College" in t or t.startswith("College")) and x > 8:
            put_text(sh, [
                "College:  IIT Gandhinagar",
                "Stream:  CSE (Dual Degree, 5-yr)",
                "Year of graduation:  2028",
            ], size=13, color=INK)
        elif sh.name == "Rectangle 23":
            put_text(sh, "CS", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        elif sh.name == "Rectangle 24":
            put_text(sh, "DH", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for sh in slide.shapes:
        if sh.name == "Rectangle 23":
            put_text(sh, "CS", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if sh.name == "Rectangle 24":
            put_text(sh, "DH", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def fill_cover(slide):
    """Keep official brand cover; stamp team + idea under the challenge line."""
    # Avoid fighting brand art — add a quiet identity strip at bottom if room.
    add_tb(
        slide, 3.2, 7.05, 6.9, 0.32,
        "ControlPlane  ·  ControlPlane.ai  ·  Accenture Innovation Challenge 2026",
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )


def fill_thank_you(slide):
    """Thank you panel — closer law, not a logo dump."""
    clear_shapes(slide)
    # dark plate matching content system
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), I(SW), I(SH))
    box.fill.solid()
    box.fill.fore_color.rgb = PLATE
    box.line.fill.background()
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), I(0.12), I(SH))
    rail.fill.solid()
    rail.fill.fore_color.rgb = RED
    rail.line.fill.background()
    add_tb(slide, 0.8, 2.4, 11.5, 0.35,
           "CONTROLPLANE.AI", size=14, bold=True, color=RED, align=PP_ALIGN.LEFT)
    add_tb(slide, 0.8, 2.9, 11.5, 1.4,
           "Now nothing acts until it can prove it should.",
           size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_tb(slide, 0.8, 4.6, 11.5, 0.5,
           "no span → no execution  ·  hold / escalate  ·  publish the miss",
           size=16, bold=False, color=C(0x8A, 0x87, 0x90), align=PP_ALIGN.LEFT)
    add_tb(slide, 0.8, 6.5, 11.5, 0.4,
           "Thank you", size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def _struct_template() -> Path:
    """Cover → Team → 3 content → Video → Thank You from the official AIC file."""
    work = Path("/tmp/aic-build")
    struct = Path("/tmp/aic-struct.pptx")
    subprocess.check_call(["rm", "-rf", str(work)])
    subprocess.check_call(["python3", str(SK / "office/unpack.py"), str(SRC), str(work)])
    # Remove instructions slide (index 2 in 1-based? delete_slide uses slide part)
    # Template: 1 cover, 2 instructions, 3 team, 4 problem, 5 solution, 6 video, 7 thanks
    # After delete instructions + add one content slide → 7 slides:
    # cover, team, content, content, content, video, thanks
    subprocess.check_call(["python3", str(SK / "delete_slide.py"), str(work), "2"])
    subprocess.check_call([
        "python3", str(SK / "add_slide.py"), str(work), "slide4.xml", "--after", "slide5.xml",
    ])
    subprocess.check_call(["python3", str(SK / "clean.py"), str(work)])
    subprocess.check_call([
        "python3", str(SK / "office/pack.py"), str(work), str(struct), "--original", str(SRC),
    ])
    return struct


def main():
    posters = {
        2: ROOT / "posters" / "s1.png",
        3: ROOT / "posters" / "s2.png",
        4: ROOT / "posters" / "s3.png",
        5: ROOT / "posters" / "sv.png",
    }
    for p in posters.values():
        assert p.exists(), f"missing poster {p}"

    prs = Presentation(str(_struct_template()))
    assert len(prs.slides) == 7, len(prs.slides)

    # 0 cover · 1 team · 2 graph · 3 matrix · 4 refuse · 5 video · 6 thank you
    fill_cover(prs.slides[0])
    fill_team(prs.slides[1])
    for idx, path in posters.items():
        full_bleed(prs.slides[idx], path)
    fill_thank_you(prs.slides[6])

    prs.save(str(OUT))
    size = _slim_pptx(OUT)
    # PDF twin of full package (for portal that wants PDF = same deck)
    _export_pdf_from_pptx(OUT)
    print(f"slides: {len(prs.slides)}")
    print(f"{OUT.name}  {size/1e6:.2f} MB  (limit 20)")
    assert size < 20 * 1024 * 1024, size
    assert len(prs.slides) == 7
    print("OK — official 7-slide package (full-bleed content)")


def _slim_pptx(path: Path) -> int:
    """Recompress large embedded JPEGs so official package stays under 20 MB."""
    import zipfile
    import io as _io

    src = zipfile.ZipFile(path, "r")
    names = list(dict.fromkeys(reversed(src.namelist())))[::-1]
    tmp = path.with_suffix(".pptx.tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as out:
        for n in names:
            d = src.read(n)
            low = n.lower()
            if low.endswith((".jpg", ".jpeg")) and len(d) > 200_000:
                try:
                    im = Image.open(_io.BytesIO(d)).convert("RGB")
                    if max(im.size) > 1920:
                        r = 1920 / max(im.size)
                        im = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                    b = _io.BytesIO()
                    im.save(b, "JPEG", quality=78, optimize=True, progressive=True)
                    if b.tell() < len(d):
                        d = b.getvalue()
                except Exception:
                    pass
            elif low.endswith(".png") and len(d) > 400_000:
                try:
                    im = Image.open(_io.BytesIO(d)).convert("RGB")
                    if max(im.size) > 1920:
                        r = 1920 / max(im.size)
                        im = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                    b = _io.BytesIO()
                    im.save(b, "JPEG", quality=86, optimize=True, progressive=True)
                    # keep as png path but JPEG bytes can break — rewrite name? skip rename; keep PNG re-encode
                    b2 = _io.BytesIO()
                    im.save(b2, "PNG", optimize=True)
                    if b2.tell() < len(d):
                        d = b2.getvalue()
                except Exception:
                    pass
            out.writestr(n, d)
    src.close()
    os.replace(tmp, path)
    return path.stat().st_size


def _export_pdf_from_pptx(pptx_path: Path):
    """Best-effort PDF of full package via posters order for content; skip if tools missing."""
    try:
        import img2pdf
        pages = []
        # Cover/team/thanks lack posters — PDF form twin is content-only via build_submission.
        # Full PDF: rasterize key pages from posters + dark thank-you still.
        order = [
            ROOT / "posters" / "s1.png",
            ROOT / "posters" / "s2.png",
            ROOT / "posters" / "s3.png",
        ]
        pdf_out = ROOT / "ControlPlane_ControlPlane-ai.pdf"
        # Keep content PDF identical to form twin for the content story;
        # portal PDF often wants the same 3 solution slides.
        jpegs = []
        for png in order:
            im = Image.open(png).convert("RGB")
            w, h = im.size
            scale = min(1.0, MAX_EDGE / max(w, h))
            if scale < 1.0:
                im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
            tmp = ROOT / "posters" / f"_fullpdf_{png.stem}.jpg"
            im.save(tmp, "JPEG", quality=88, optimize=True)
            jpegs.append(str(tmp))
        layout = img2pdf.get_layout_fun(
            pagesize=(img2pdf.in_to_pt(SW), img2pdf.in_to_pt(SH))
        )
        pdf_out.write_bytes(img2pdf.convert(jpegs, layout_fun=layout))
        for j in jpegs:
            try:
                os.remove(j)
            except OSError:
                pass
        print(f"{pdf_out.name}  {pdf_out.stat().st_size/1e6:.2f} MB  (3 content pages)")
    except Exception as e:
        print("PDF export skipped:", e)


if __name__ == "__main__":
    main()
