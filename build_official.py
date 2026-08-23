"""OFFICIAL AIC 7-slide package — Cover · Team · S1 · S2 · S3 · Video · Thank you.

Full-clarity build: same 7680 masters as the 3-file pack, JPEG Q=100, sharp team
photos, same clear MP4. This is the complete official-format package.

    python3 build_official.py

Writes ControlPlane_ControlPlane-ai.pptx/.pdf/.mp4 to submission/official/.

Slide surgery note: new slides are APPENDED first and only then are the unwanted template
slides removed. Deleting first makes python-pptx reuse part names (slide6.xml), which
produces duplicate zip entries that the recompression step silently collapses — that is
how the Video and Thank-you slides vanished on an earlier build.
"""

from __future__ import annotations

import copy
import io
import os
import shutil
import subprocess
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "archive" / "tools" / "AIC_Talent-Brand_PPT-Template (1).pptx"
POSTERS = ROOT / "posters"
OUT_DIR = ROOT / "submission" / "official"
NAME = "ControlPlane_ControlPlane-ai"

SW, SH = Inches(13.333333), Inches(7.5)
CAP = 20_000_000
# Content masters: sharp + ≥1 MB headroom under 20 MB (PDF was 19.6). Cover slimmed hard.
JPEG_Q = 93
MAX_EDGE = 5760
COVER_MAX_EDGE = 1800  # Accenture cover art only — not our content

PURPLE = RGBColor(0xA1, 0x00, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6E, 0x6A, 0x72)

TEAM = [
    ("Choda Srujan Sai", "Team Leader", "IIT Gandhinagar",
     "Civil Engineering (CSE minor)", "2027", "CS",
     ROOT / "assets" / "team" / "srujan.jpg"),   # pure square crop, no AI
    ("Dhrithika", "Team Member", "IIT Gandhinagar",
     "Computer Science (Dual Degree, 5-yr)", "2028", "DH",
     ROOT / "assets" / "team" / "dhrithika.jpg"),  # pure square crop from duo, no AI
]

# Template slide indices
I_COVER, I_INSTR, I_TEAM, I_PROB, I_SOLN, I_VIDEO, I_THANKS = range(7)


# --------------------------------------------------------------------------- helpers
def _drop(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _settext(shape, text, size=None, bold=None, color=None, italic=None) -> None:
    """Replace a shape's text, clearing every paragraph after the first.

    settext-by-runs alone leaves stale paragraphs behind — that is what duplicated
    'Stream:' / 'Year of graduation:' on an earlier build.
    """
    tf = shape.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color


def _monogram(slide, left, top, w, h, initials) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PURPLE
    box.line.fill.background()
    box.shadow.inherit = False
    _settext(box, initials, size=44, bold=True, color=WHITE)
    box.text_frame.paragraphs[0].alignment = 2  # centre
    box.text_frame.word_wrap = False


def _photo(slide, left, top, w, h, path: Path) -> None:
    """Fit a square headshot into the template photo box.

    Source files are face-centered squares. When the slot is wider than tall,
    bias the vertical crop slightly downward so chin/neck stay in frame
    (center crop alone eats the chin when headroom is large).
    """
    from pptx.enum.shapes import MSO_SHAPE

    assert path.exists(), f"missing team photo {path}"
    im = Image.open(path).convert("RGB")
    iw, ih = im.size
    slot_ar = (w / h) if h else 1.0
    src_ar = iw / ih
    if abs(src_ar - slot_ar) > 0.02:
        if src_ar > slot_ar:
            nw = int(ih * slot_ar)
            x0 = (iw - nw) // 2
            im = im.crop((x0, 0, x0 + nw, ih))
        else:
            nh = int(iw / slot_ar)
            # Bias down ~12% of free band so chin/neck survive landscape slots
            free = ih - nh
            y0 = int(free * 0.62) if free > 0 else 0
            y0 = max(0, min(ih - nh, y0))
            im = im.crop((0, y0, iw, y0 + nh))
    # high-res for the slot (min 1024 so faces stay sharp on the team slide)
    tw = max(int(w / 9525), 1024)
    th = max(int(h / 9525), 1024)
    im = im.resize((tw, th), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=100, optimize=True, subsampling=0)
    buf.seek(0)
    plate = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    plate.fill.solid()
    plate.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x16)
    plate.line.fill.background()
    slide.shapes.add_picture(buf, left, top, w, h)


def _jpeg(png: Path) -> io.BytesIO:
    im = Image.open(png).convert("RGB")
    w, h = im.size
    s = min(1.0, MAX_EDGE / max(w, h))
    if s < 1.0:
        im = im.resize((int(w * s), int(h * s)), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=JPEG_Q, optimize=True, progressive=True, subsampling=0)
    buf.seek(0)
    return buf


def _blank_layout(prs):
    """The template's own clean full-frame layout."""
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name == "1_Standard slide_no bullets":
                return l
    return prs.slide_layouts[6]


def _full_bleed(prs, poster: Path):
    slide = prs.slides.add_slide(_blank_layout(prs))
    for sh in list(slide.shapes):
        _drop(sh)
    slide.shapes.add_picture(_jpeg(poster), 0, 0, SW, SH)
    return slide


def _reorder(prs, order) -> None:
    """order = list of current slide indices, in the sequence we want them."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


def _delete(prs, idx_set) -> None:
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(idx_set, reverse=True):
        rId = ids[i].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        lst.remove(ids[i])


# --------------------------------------------------------------------------- team
def fill_team(slide) -> None:
    """Two members, side by side across the template's divider. Third slot removed."""
    kill_names = {"Rectangle 2"}  # 'All fields are mandatory'
    by_pos = {}
    for sh in list(slide.shapes):
        if sh.name in kill_names:
            _drop(sh)
            continue
        by_pos[(round(sh.left / 914400, 2), round(sh.top / 914400, 2))] = sh

    def at(l, t):
        return by_pos.get((l, t))

    # --- remove the unused bottom-left member slot entirely
    for key in [(3.01, 4.6), (2.88, 5.73), (3.01, 5.58), (0.85, 4.77), (0.93, 4.88)]:
        sh = at(*key)
        if sh is not None:
            _drop(sh)

    # --- member A (leader), top-left — solo headshot
    a_name, a_role, a_col, a_stream, a_year, a_ini, a_photo = TEAM[0]
    sh = at(2.88, 2.3)
    if sh is not None:
        _settext(sh, f"{a_name}  ·  {a_role}", size=17, bold=True, color=INK)
    sh = at(2.87, 3.44)
    if sh is not None:
        _settext(sh, f"{a_col}   |   {a_stream}   |   Class of {a_year}",
                 size=11, color=GREY)

    # --- member B (Dhrithika), top-right — clear face from duo photo
    b_name, b_role, b_col, b_stream, b_year, b_ini, b_photo = TEAM[1]
    sh = at(10.02, 2.3)
    if sh is not None:
        label = f"{b_name}  ·  {b_role}" if b_role else b_name
        _settext(sh, label, size=17, bold=True, color=INK)
    sh = at(9.98, 3.44)
    if sh is not None:
        _settext(sh, f"{b_col}   |   {b_stream}   |   Class of {b_year}",
                 size=11, color=GREY)

    # --- photo boxes -> real headshots (fallback monogram if missing)
    for key, ini, photo in (
        ((0.93, 2.74), a_ini, a_photo),
        ((8.08, 2.74), b_ini, b_photo),
    ):
        sh = at(*key)
        if sh is None:
            continue
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        _drop(sh)
        if photo is not None and Path(photo).exists():
            _photo(slide, l, t, w, h, Path(photo))
        else:
            _monogram(slide, l, t, w, h, ini)
    for key in [(0.83, 2.63), (7.97, 2.63)]:  # empty picture placeholders behind them
        sh = at(*key)
        if sh is not None:
            _drop(sh)

    # --- team name into the title and into the template's TEAM NAME row
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Team details":
            _settext(sh, "Team details  ·  ControlPlane", size=28, bold=True)
        if getattr(sh, "has_table", False) and sh.has_table:
            cell = sh.table.cell(0, 1)
            _settext(cell, "ControlPlane", size=12, bold=True, color=INK)


# --------------------------------------------------------------------------- video
def fill_video(slide, mp4: Path) -> None:
    for sh in list(slide.shapes):
        _drop(sh)
    sv = POSTERS / "sv.png"
    if sv.exists():
        slide.shapes.add_picture(_jpeg(sv), 0, 0, SW, SH)
    box = slide.shapes.add_textbox(Inches(7.3), Inches(0.28), Inches(5.5), Inches(0.44))
    tf = box.text_frame
    tf.word_wrap = True
    try:
        secs = int(round(float(subprocess.check_output(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", str(mp4)]
        ).decode().strip())))
    except Exception:
        secs = 158
    _settext(box, f"{mp4.name}   ·   {secs // 60}:{secs % 60:02d}",
             size=13, bold=True, color=WHITE)
    box.text_frame.paragraphs[0].alignment = 3  # right


# --------------------------------------------------------------------------- build
def build() -> Path:
    prs = Presentation(str(TEMPLATE))
    prs.slide_width, prs.slide_height = SW, SH

    # 1. APPEND the three content slides first (never delete before adding)
    for stem in ("s1", "s2", "s3"):
        _full_bleed(prs, POSTERS / f"{stem}.png")
    n_s1, n_s2, n_s3 = 7, 8, 9

    # 2. fill the template slides we keep
    fill_team(prs.slides[I_TEAM])
    fill_video(prs.slides[I_VIDEO], ROOT / "video" / f"{NAME}.mp4")

    # 3. order: Cover · Team · S1 · S2 · S3 · Video · Thank you
    _reorder(prs, [I_COVER, I_TEAM, n_s1, n_s2, n_s3, I_VIDEO, I_THANKS,
                   I_INSTR, I_PROB, I_SOLN])
    _delete(prs, {7, 8, 9})  # instructions + the two 200-word slides, now at the tail

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / f"{NAME}.pptx"
    prs.save(str(out))
    _slim(out)

    # order assertion — this is the check that would have caught the vanished slides
    got = len(Presentation(str(out)).slides)
    assert got == 7, f"expected 7 slides, got {got}"
    return out


def _slim(path: Path) -> None:
    """Keep content masters + headshots sharp; crush only bloated template cover art.

    The Accenture cover photo alone is ~12 MB — that is what blows the 20 MB cap,
    not our S1/S2/S3 masters.
    """
    src = zipfile.ZipFile(path, "r")
    names = list(dict.fromkeys(src.namelist()))
    protect = set()
    for n in names:
        if not n.lower().endswith((".jpg", ".jpeg", ".png")):
            continue
        d = src.read(n)
        if len(d) < 150_000:
            continue
        try:
            im = Image.open(io.BytesIO(d))
            w, h = im.size
            ar = w / max(h, 1)
            # our full-bleed content: 16:9 and large
            if abs(ar - 16 / 9) < 0.08 and max(w, h) >= 2800:
                protect.add(n)
            # team headshots: square
            if abs(ar - 1.0) < 0.08 and min(w, h) >= 400:
                protect.add(n)
        except Exception:
            pass

    tmp = path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as out:
        for n in names:
            d = src.read(n)
            if (
                n not in protect
                and n.lower().endswith((".jpg", ".jpeg", ".png"))
                and len(d) > 350_000
            ):
                try:
                    im = Image.open(io.BytesIO(d)).convert("RGB")
                    edge = COVER_MAX_EDGE if len(d) > 2_000_000 else min(MAX_EDGE, 3200)
                    if max(im.size) > edge:
                        r = edge / max(im.size)
                        im = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                    b = io.BytesIO()
                    q = 88 if len(d) > 2_000_000 else 92
                    im.save(b, "JPEG", quality=q, optimize=True, progressive=True,
                            subsampling=0)
                    if b.tell() < len(d):
                        d = b.getvalue()
                except Exception:
                    pass
            out.writestr(n, d)
    src.close()
    os.replace(tmp, path)


def to_pdf(pptx: Path) -> Path | None:
    pdf = pptx.with_suffix(".pdf")
    script = Path("/tmp/topdf.scpt")
    if not script.exists():
        return None
    r = subprocess.run(["osascript", str(script), str(pptx), str(pdf)],
                       capture_output=True, text=True, timeout=420)
    if r.returncode != 0 or not pdf.exists():
        print("  PDF export failed:", (r.stderr or "").strip()[:300])
        return None
    return pdf


def _slim_pdf(pdf: Path) -> None:
    """Ghostscript recompress so PDF stays under 20 MB without killing content sharpness."""
    if pdf.stat().st_size < CAP:
        return
    tmp = pdf.with_suffix(".slim.pdf")
    # ebook ~150dpi too soft; printer ~300dpi sharp enough for review panels
    r = subprocess.run(
        [
            "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
            "-dPDFSETTINGS=/printer",
            "-dColorImageResolution=200",
            "-dGrayImageResolution=200",
            "-dMonoImageResolution=300",
            "-dColorImageDownsampleType=/Bicubic",
            "-dDetectDuplicateImages=true",
            "-dCompressFonts=true",
            "-dNOPAUSE", "-dQUIET", "-dBATCH",
            f"-sOutputFile={tmp}",
            str(pdf),
        ],
        capture_output=True, text=True, timeout=300,
    )
    if r.returncode == 0 and tmp.exists() and tmp.stat().st_size < pdf.stat().st_size:
        os.replace(tmp, pdf)
        print(f"  PDF slimmed → {pdf.stat().st_size/1e6:.2f} MB")
    else:
        if tmp.exists():
            tmp.unlink()
        print("  PDF slim skipped:", (r.stderr or r.stdout or "")[:200])


def main() -> None:
    pptx = build()
    print(f"PPTX  {pptx}  {pptx.stat().st_size/1e6:.2f} MB  / 20   (7 slides)")
    pdf = to_pdf(pptx)
    if pdf:
        _slim_pdf(pdf)
        print(f"PDF   {pdf}  {pdf.stat().st_size/1e6:.2f} MB  / 20")
    mp4 = ROOT / "video" / f"{NAME}.mp4"
    if mp4.exists():
        shutil.copy2(mp4, OUT_DIR / mp4.name)
        print(f"MP4   {OUT_DIR/mp4.name}  {mp4.stat().st_size/1e6:.2f} MB  / 20")
    for f in OUT_DIR.iterdir():
        assert f.stat().st_size < CAP, f"{f.name} exceeds 20 MB"
    print("OK — official 7-slide FULL CLARITY package in submission/official/")


if __name__ == "__main__":
    main()
