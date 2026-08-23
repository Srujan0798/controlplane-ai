#!/usr/bin/env python3
"""FORM UPLOAD — exactly 3 files in submission/. Quality from resolution, not grain air.

Hard limits: PPTX ≤20 MB · PDF ≤20 MB · MP4 ≤20 MB.
Masters are 7680×4320 (dpr=4). Embed as JPEG Q=100 / 4:4:4 — real type edges,
not seeded noise padding.
"""
from __future__ import annotations

import os
import shutil
from io import BytesIO
from pathlib import Path

import img2pdf
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

Image.MAX_IMAGE_PIXELS = None

ROOT = Path(__file__).resolve().parent
SUB = ROOT / "submission"
OUT_PPTX = SUB / "ControlPlane_ControlPlane-ai.pptx"
OUT_PDF = SUB / "ControlPlane_ControlPlane-ai.pdf"
OUT_MP4 = SUB / "ControlPlane_ControlPlane-ai.mp4"

POSTERS = [
    ROOT / "posters" / "s1.png",
    ROOT / "posters" / "s2.png",
    ROOT / "posters" / "s3.png",
]
SW, SH = 13.333333, 7.5
MAX_EDGE = 7680
JPEG_Q = 100


def _fit(im: Image.Image) -> Image.Image:
    w, h = im.size
    scale = min(1.0, MAX_EDGE / max(w, h))
    if scale < 1.0:
        im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    return im


def _to_jpeg(path: Path) -> BytesIO:
    im = _fit(Image.open(path).convert("RGB"))
    buf = BytesIO()
    # Q=100, no chroma sub-sample — hairlines and type edges survive.
    im.save(buf, "JPEG", quality=JPEG_Q, optimize=True, progressive=False, subsampling=0)
    buf.seek(0)
    return buf


def build_pptx() -> int:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]
    for png in POSTERS:
        assert png.exists(), f"missing {png} — run python3 -m visuals.render"
        slide = prs.slides.add_slide(blank)
        for sh in list(slide.shapes):
            sp = sh._element
            sp.getparent().remove(sp)
        slide.shapes.add_picture(
            _to_jpeg(png), Emu(0), Emu(0), prs.slide_width, prs.slide_height
        )
    assert len(prs.slides) == 3
    SUB.mkdir(exist_ok=True)
    prs.save(OUT_PPTX)
    return OUT_PPTX.stat().st_size


def build_pdf() -> int:
    pages = []
    tmp_paths = []
    for png in POSTERS:
        tmp = ROOT / "posters" / f"_pdf_{png.stem}.jpg"
        buf = _to_jpeg(png)
        tmp.write_bytes(buf.getvalue())
        pages.append(str(tmp))
        tmp_paths.append(tmp)
    layout = img2pdf.get_layout_fun(pagesize=(img2pdf.in_to_pt(SW), img2pdf.in_to_pt(SH)))
    data = img2pdf.convert(pages, layout_fun=layout)
    SUB.mkdir(exist_ok=True)
    OUT_PDF.write_bytes(data)
    for p in tmp_paths:
        try:
            os.remove(p)
        except OSError:
            pass
    return OUT_PDF.stat().st_size


def sync_video() -> None:
    src = ROOT / "video" / "ControlPlane_ControlPlane-ai.mp4"
    assert src.exists() and src.stat().st_size > 100_000, "missing video — run build_video.py"
    shutil.copy2(src, OUT_MP4)
    print(f"video: {OUT_MP4}  {OUT_MP4.stat().st_size/(1024*1024):.2f} MB  / 20")


def purge_non_upload() -> None:
    SUB.mkdir(exist_ok=True)
    keep = {OUT_PPTX.name, OUT_PDF.name, OUT_MP4.name}
    for p in list(SUB.iterdir()):
        if p.name in keep:
            continue
        if p.is_dir():
            shutil.rmtree(p)
            print(f"removed dir {p.relative_to(ROOT)}")
        else:
            p.unlink()
            print(f"removed {p.relative_to(ROOT)}")


def main():
    for p in POSTERS:
        im = Image.open(p)
        print(f"poster {p.name}: {im.size[0]}x{im.size[1]}")
        assert max(im.size) >= 5760, f"{p} too small — re-render at dpr=4"

    n = build_pptx()
    m = build_pdf()
    assert n < 20 * 1024 * 1024 and m < 20 * 1024 * 1024, "over 20 MB — lower Q"
    sync_video()
    purge_non_upload()

    only = sorted(p.name for p in SUB.iterdir())
    assert only == sorted([OUT_PPTX.name, OUT_PDF.name, OUT_MP4.name]), only

    print(f"slides: 3")
    print(f"PPTX  {OUT_PPTX}  {n/(1024*1024):.2f} MB  / 20  (JPEG Q={JPEG_Q}, no grain pad)")
    print(f"PDF   {OUT_PDF}   {m/(1024*1024):.2f} MB  / 20  (JPEG Q={JPEG_Q}, no grain pad)")
    print(f"MP4   {OUT_MP4}   {OUT_MP4.stat().st_size/(1024*1024):.2f} MB  / 20")
    print("OK — form upload ready: exactly 3 files in submission/")


if __name__ == "__main__":
    main()
